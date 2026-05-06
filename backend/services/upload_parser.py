"""Module for parsing uploaded document files (PDF, DOCX)."""
import io
import os
from typing import Dict, Any, Optional, Tuple, List
from pptx import Presentation
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document as DocxDocument

def _parse_pdf(path: str) -> Tuple[Optional[str], Optional[int]]:
    with open(path, "rb") as fh:
        data = fh.read()
    with fitz.open(stream=data, filetype="pdf") as doc:
        pages = doc.page_count
        out = []
        for i in range(pages):
            out.append(doc.load_page(i).get_text("text"))
        txt = "\n".join(out).strip()
        return (txt if txt else None), pages

def _parse_docx(path: str) -> Tuple[Optional[str], Optional[int]]:
    with open(path, "rb") as fh:
        bio = io.BytesIO(fh.read())
    doc = DocxDocument(bio)
    txt = "\n".join([p.text for p in doc.paragraphs]).strip()
    return (txt if txt else None), None

def _parse_txt(path: str) -> Tuple[Optional[str], Optional[int]]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(), None
    except Exception:
        return None, None


def _parse_pptx(path: str) -> Tuple[Optional[str], Optional[int]]:
    try:
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)
        return "\n".join(out).strip(), len(prs.slides)
    except Exception:
        return None, None


def parse_document(path: str) -> Dict[str, Any]:
    ext = os.path.splitext(path)[1].lower()
    out: Dict[str, Any] = {"ext": ext.lstrip(".")}
    try:
        if ext == ".pdf":
            text, pages = _parse_pdf(path)
            out["text"] = text
            out["pages"] = pages
        elif ext == ".docx":
            text, _ = _parse_docx(path)
            out["text"] = text
        elif ext == ".pptx":
            text, slides = _parse_pptx(path)
            out["text"] = text
            out["slides"] = slides
        elif ext in [".txt", ".md"]:
            text, _ = _parse_txt(path)
            out["text"] = text
        else:
            pass
    except Exception as e:
        out["error"] = str(e)
    return out

def extract_text_from_file(path: str) -> str:
    p = Path(path)
    ext = p.suffix.lower()

    if ext == ".txt":
        return p.read_text(errors="ignore")

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(p))
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)
        return "\n".join(out)

    if ext == ".docx":
        import docx
        d = docx.Document(str(p))
        return "\n".join([para.text for para in d.paragraphs])

    if ext == ".pdf":
        # minimal PDF extraction (works for text-based PDFs)
        from pypdf import PdfReader
        reader = PdfReader(str(p))
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return "\n".join(parts)

    return ""